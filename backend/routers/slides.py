"""
Slides-to-video router.
Accepts a .pptx file, converts each slide to a PNG using LibreOffice,
then stitches them into a silent .mp4 using FFmpeg with optional transitions.
Frames are stored persistently so the merge step can rebuild with custom timing.
"""
import os
import glob
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Optional, List

from fastapi import APIRouter, UploadFile, File, Form, BackgroundTasks, HTTPException

from config import settings
from job_store import job_store

router = APIRouter()


def find_libreoffice() -> str:
    """Find the LibreOffice binary on the current system."""
    if settings.libreoffice_binary not in ("libreoffice", "soffice"):
        if os.path.exists(settings.libreoffice_binary):
            return settings.libreoffice_binary

    # Fixed known Windows paths
    windows_paths = [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]
    # Also glob for any versioned install (LibreOffice 7, 24, 24.2, 25.x, etc.)
    for base in [r"C:\Program Files", r"C:\Program Files (x86)"]:
        windows_paths += glob.glob(os.path.join(base, "LibreOffice*", "program", "soffice.exe"))

    for path in windows_paths:
        if os.path.exists(path):
            return path

    # Check system PATH as a last resort
    which_soffice = shutil.which("soffice")
    if which_soffice and os.path.exists(which_soffice):
        return which_soffice
    which_lo = shutil.which("libreoffice")
    if which_lo and os.path.exists(which_lo):
        return which_lo

    raise RuntimeError(
        "LibreOffice is not installed on this computer.\n\n"
        "To use the PPTX converter, please:\n"
        "1. Download LibreOffice from https://www.libreoffice.org/download/\n"
        "2. Install it (keep the default installation path)\n"
        "3. Restart ScriptToVideo"
    )


def convert_pptx_to_images(pptx_path: str, output_dir: str) -> list:
    """
    Convert PPTX to images via two steps:
      1. LibreOffice: PPTX -> PDF  (all slides reliably exported)
      2. PyMuPDF (fitz): PDF pages -> PNG images
    """
    import fitz  # PyMuPDF

    libreoffice_bin = find_libreoffice()

    # Use a dedicated per-conversion user profile so headless LibreOffice never
    # conflicts with a running GUI instance (which causes [WinError 2] / lock errors).
    lo_profile_dir = os.path.join(output_dir, "_lo_profile")
    os.makedirs(lo_profile_dir, exist_ok=True)
    # LibreOffice expects a file:// URI for UserInstallation
    lo_profile_uri = "file:///" + lo_profile_dir.replace("\\", "/").lstrip("/")

    cmd = [
        libreoffice_bin,
        "--headless",
        "--norestore",           # don't try to restore a previous session
        "--nofirststartwizard",  # skip first-run wizard
        f"-env:UserInstallation={lo_profile_uri}",  # isolated profile per job
        "--convert-to", "pdf",
        "--outdir", output_dir,
        pptx_path,
    ]
    try:
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=600)  # 10 min for large decks
    except FileNotFoundError:
        raise RuntimeError(
            "LibreOffice is not installed on this computer.\n\n"
            "To use the PPTX converter, please:\n"
            "1. Download LibreOffice from https://www.libreoffice.org/download/\n"
            "2. Install it (keep the default installation path)\n"
            "3. Restart ScriptToVideo"
        )
    except subprocess.TimeoutExpired:
        raise RuntimeError(
            "LibreOffice took longer than 10 minutes to convert the presentation. "
            "Try splitting the file into smaller parts (50-60 slides each)."
        )
    if result.returncode != 0:
        stderr = (result.stderr or "").strip()
        raise RuntimeError(f"LibreOffice PDF export failed (code {result.returncode}): {stderr}")

    pdf_files = glob.glob(os.path.join(output_dir, "*.pdf"))
    if not pdf_files:
        raise RuntimeError("LibreOffice did not produce a PDF file.")
    pdf_path = pdf_files[0]

    doc = fitz.open(pdf_path)
    images = []
    for i, page in enumerate(doc):
        mat = fitz.Matrix(2.0, 2.0)
        pix = page.get_pixmap(matrix=mat)
        img_path = os.path.join(output_dir, f"slide_{i+1:03d}.png")
        pix.save(img_path)
        images.append(img_path)
    doc.close()

    # Remove the intermediate PDF (keep only PNGs)
    try:
        os.remove(pdf_path)
    except Exception:
        pass

    return sorted(images)


def _iter_shapes(shapes):
    """
    Recursively yield all leaf shapes, descending into group shapes.
    python-pptx's slide.shapes only returns top-level shapes — grouped
    shapes must be walked manually.
    """
    for shape in shapes:
        try:
            # GroupShape has a .shapes attribute
            if shape.shape_type == 6:   # MSO_SHAPE_TYPE.GROUP == 6
                yield from _iter_shapes(shape.shapes)
                continue
        except Exception:
            pass
        yield shape


def extract_slide_texts(pptx_path: str, n_slides: int) -> dict:
    """
    Extract title and body text from each slide using python-pptx.
    Returns {"titles": [...], "bodies": [...]} — titles are used for
    exact matching in AI sync since they are identical in the voiceover.

    Title detection strategy (most-specific first):
      1. Shape with PP_PLACEHOLDER.TITLE / CENTER_TITLE / VERTICAL_TITLE
      2. Shape whose placeholder idx == 0 (index-0 placeholder is always the title)
      3. Shape whose name contains "title" (case-insensitive)
      4. First text shape on the slide (fallback — better than nothing)

    Group shapes are recursively walked so text inside groups is found.
    """
    try:
        from pptx import Presentation
        from pptx.enum.shapes import PP_PLACEHOLDER

        TITLE_PH_TYPES = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE,
                          PP_PLACEHOLDER.VERTICAL_TITLE}

        prs = Presentation(pptx_path)
        titles, bodies = [], []

        for slide in prs.slides:
            title_text = ""
            body_parts = []
            first_text_shape = None   # fallback

            title_candidate = None
            title_priority  = 99      # lower = better match

            # Walk ALL shapes including those nested inside groups
            for shape in _iter_shapes(slide.shapes):
                if not shape.has_text_frame:
                    continue
                try:
                    text = shape.text_frame.text.strip()
                except Exception:
                    continue
                if not text:
                    continue

                if first_text_shape is None:
                    first_text_shape = text

                try:
                    ph = shape.placeholder_format
                except (ValueError, AttributeError):
                    ph = None
                ph_type  = ph.type if ph else None
                ph_idx   = ph.idx  if ph else None

                priority = 99
                if ph_type in TITLE_PH_TYPES:
                    priority = 0
                elif ph_idx == 0:
                    priority = 1
                elif shape.name and "title" in shape.name.lower():
                    priority = 2

                if priority < title_priority:
                    title_candidate = text
                    title_priority  = priority

            # Assign title — if no placeholder found, use the TOPMOST text
            # shape (lowest Y value) rather than the first in XML order.
            # Slides like "HOW AED WORKS" have their header shape listed late
            # in the XML but positioned at the top of the slide visually.
            # We skip all-emoji/symbol shapes (< 3 ASCII letters) so that
            # decorative icon shapes (🚨, ⚡, etc.) at Y≈0 aren't picked.
            if title_candidate:
                title_text = title_candidate
            else:
                import re as _re
                def _has_word_chars(s: str) -> bool:
                    return len(_re.sub(r"[^A-Za-z0-9]", "", s)) >= 3

                topmost_text = None
                topmost_y = float("inf")
                for shape in _iter_shapes(slide.shapes):
                    if not shape.has_text_frame:
                        continue
                    try:
                        text = shape.text_frame.text.strip()
                    except Exception:
                        continue
                    if not text or not _has_word_chars(text):
                        continue
                    y = getattr(shape, "top", None) or 0
                    if y < topmost_y:
                        topmost_y = y
                        topmost_text = text
                title_text = topmost_text or first_text_shape or ""

            # Collect body text (all non-title shapes)
            for shape in _iter_shapes(slide.shapes):
                if not shape.has_text_frame:
                    continue
                try:
                    text = shape.text_frame.text.strip()
                except Exception:
                    continue
                if not text or text == title_text:
                    continue
                for para in shape.text_frame.paragraphs:
                    try:
                        line = " ".join(run.text for run in para.runs).strip()
                    except Exception:
                        line = para.text.strip()
                    if line:
                        body_parts.append(line)

            titles.append(title_text)
            bodies.append(" ".join(body_parts))

        while len(titles) < n_slides:
            titles.append("")
            bodies.append("")

        return {"titles": titles[:n_slides], "bodies": bodies[:n_slides]}
    except Exception as exc:
        import traceback
        traceback.print_exc()
        print(f"[extract_slide_texts] WARNING: fell back to empty titles — {exc}")
        return {"titles": [""] * n_slides, "bodies": [""] * n_slides}


# ── Text-animation helpers ────────────────────────────────────────────────────

def get_element_boxes(pptx_path: str, slide_idx: int, img_w: int, img_h: int) -> list:
    """
    Return a list of (l, t, r, b, kind) tuples for every animatable element on
    a slide, sorted top→bottom by their vertical midpoint.

    kind values: "text" | "image"

    Text shapes are split into per-paragraph rows so each line animates
    individually.  The paragraph heights are approximated by dividing the
    shape's pixel height equally among non-empty paragraphs — not pixel-perfect
    but visually convincing.

    Image shapes (pictures / filled shapes without text) are returned as-is.
    """
    try:
        from pptx import Presentation
        from pptx.util import Pt
        prs = Presentation(pptx_path)
        if slide_idx >= len(prs.slides):
            return []
        slide = prs.slides[slide_idx]
        sx = img_w / prs.slide_width
        sy = img_h / prs.slide_height
        elements = []

        for shape in slide.shapes:
            sl = max(0, int(shape.left  * sx))
            st = max(0, int(shape.top   * sy))
            sr = min(img_w, int((shape.left + shape.width)  * sx))
            sb = min(img_h, int((shape.top  + shape.height) * sy))

            # ── LINE shapes are inherently thin — give them a minimum box ────
            _is_line = (getattr(shape, "shape_type", None) == 10)
            if _is_line:
                if sr <= sl: sr = min(img_w, sl + 4)
                if sb <= st: sb = min(img_h, st + 4)
            elif sr <= sl + 10 or sb <= st + 10:
                continue

            _stype = getattr(shape, "shape_type", None)

            # ── Text shape ───────────────────────────────────────────────────
            if shape.has_text_frame:
                tf = shape.text_frame
                # Collect non-empty paragraph texts
                paras = [p for p in tf.paragraphs
                         if "".join(r.text for r in p.runs).strip()]
                if not paras:
                    continue

                # AUTO_SHAPE / CALLOUT / FREEFORM with a text frame are custom
                # coloured boxes (e.g. "with defined roles" callout cards).
                # They have their OWN fill colour distinct from the slide BG.
                # If we split them into text rows and erase each row separately
                # we leave the coloured box background visible at frame 0 while
                # the text flickers in — ugly artifact.
                # Fix: treat the whole box as a single "image" unit so the
                # background and text reveal together as one block.
                # Standard title / body placeholders (is_placeholder=True) are
                # NOT coloured boxes and must still be split into rows.
                _is_autoshape_box = (
                    _stype in (1, 2, 5) and          # AUTO_SHAPE / CALLOUT / FREEFORM
                    not getattr(shape, "is_placeholder", False)
                )
                if _is_autoshape_box:
                    shape_area = max(1, sr - sl) * max(1, sb - st)
                    if shape_area < img_w * img_h * 0.80:   # skip full-bleed BG rects
                        elements.append((sl, st, sr, sb, "image"))
                    continue

                # Regular placeholder / text box: split into per-paragraph rows
                n_p = len(paras)
                shape_h = sb - st
                row_h   = max(4, shape_h // n_p)
                for pi, _ in enumerate(paras):
                    pt = st + pi * row_h
                    # Last row extends all the way to the shape bottom so no
                    # fraction of later-paragraph text is ever left un-erased.
                    pb = sb if pi == n_p - 1 else min(sb, st + (pi + 1) * row_h)
                    if pb > pt + 2:
                        elements.append((sl, pt, sr, pb, "text"))
                continue

            # ── Table shape: split into per-row animatable bands ─────────────
            # shape_type 19 = MSO_SHAPE_TYPE.TABLE
            if _stype == 19 or getattr(shape, "has_table", False):
                try:
                    tbl = shape.table
                    n_rows = len(tbl.rows)
                    if n_rows > 0:
                        row_h = max(4, (sb - st) // n_rows)
                        for ri in range(n_rows):
                            rt = st + ri * row_h
                            rb = min(sb, st + (ri + 1) * row_h)
                            if rb > rt + 2:
                                elements.append((sl, rt, sr, rb, "table_row"))
                except Exception:
                    elements.append((sl, st, sr, sb, "image"))
                continue

            # ── Image / picture shape ────────────────────────────────────────
            # shape_type 13 = MSO_SHAPE_TYPE.PICTURE
            try:
                from pptx.enum.shapes import MSO_SHAPE_TYPE
                is_pic = shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            except Exception:
                is_pic = (_stype == 13)

            if is_pic:
                # NEVER erase picture/image shapes.  Images are reference
                # material that should be visible from the first frame.
                # Full-screen images get Ken Burns (separate path); smaller
                # inset images should simply remain static on screen.
                # Erasing images then trying to reveal them via element_timing
                # produces blank slides whenever the delay exceeds the slide
                # duration (which happens often due to timing drift).
                continue

            # ── Group shapes (type 6) ────────────────────────────────────────
            # GROUP shapes contain nested child shapes (e.g. image placeholder
            # + caption, icon clusters).  Without handling them here the group
            # bounding box is never erased and shows at frame 0 as a coloured
            # block (the teal box on the "Log Rolling" title slide, for example).
            if _stype == 6:
                shape_area = max(1, sr - sl) * max(1, sb - st)
                if shape_area < img_w * img_h * 0.80:
                    elements.append((sl, st, sr, sb, "image"))
                continue

            # ── Line / connector shapes ─────────────────────────────────────
            # MSO_SHAPE_TYPE.LINE = 10; these are animatable structural elements
            # (colored rule lines, dividers, connectors).  We treat them with
            # kind "line" so the erase pass can use a larger padding to cover
            # arrowheads / endpoints that extend outside the bounding box.
            if _is_line:
                elements.append((sl, st, sr, sb, "line"))
                continue

        # Sort by vertical midpoint so reading order is respected
        elements.sort(key=lambda e: (e[1] + e[3]) / 2)
        return elements
    except Exception:
        import traceback
        traceback.print_exc()
        return []


# Keep old name as alias so any other callers don't break
def get_text_boxes(pptx_path: str, slide_idx: int, img_w: int, img_h: int) -> list:
    return [(l, t, r, b) for l, t, r, b, _ in get_element_boxes(pptx_path, slide_idx, img_w, img_h)]


def get_slide_paragraph_texts(pptx_path: str, slide_idx: int) -> list:
    """
    Return the text of every animatable element on a slide in the same
    top-to-bottom order as get_element_boxes().

    Used by compute_element_timestamps() in sync.py to match each visual
    element to the moment it is spoken in the voiceover.

    CRITICAL: this function MUST return exactly the same number of elements
    in the same order as get_element_boxes() so that element_timing delays
    line up 1-to-1 with animation targets.  Any shape that get_element_boxes
    skips or adds must be treated identically here.

    • Text paragraphs  → paragraph text string
    • Table rows       → space-joined cell text
    • Autoshape boxes  → combined text (one entry, matching get_element_boxes)
    • Group shapes     → empty string  (only those with area < 80 % of slide)
    • Line connectors  → empty string
    • Pictures         → SKIPPED (get_element_boxes never adds them)
    """
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        if slide_idx >= len(prs.slides):
            return []
        slide      = prs.slides[slide_idx]
        slide_area = prs.slide_width * prs.slide_height  # EMU²
        items      = []   # [(raw_mid_y, text)]

        for shape in slide.shapes:
            l_r = shape.left  or 0
            t_r = shape.top   or 0
            r_r = l_r + (shape.width  or 0)
            b_r = t_r + (shape.height or 0)

            _stype_r = getattr(shape, "shape_type", None)
            _is_line_r = (_stype_r == 10)

            # ── Mirror the thin-shape guard in get_element_boxes ─────────────
            if _is_line_r:
                if r_r <= l_r: r_r = l_r + 10
                if b_r <= t_r: b_r = t_r + 10
            elif (r_r - l_r) < 10 or (b_r - t_r) < 10:
                continue

            # ── Text shapes ───────────────────────────────────────────────────
            if shape.has_text_frame:
                tf    = shape.text_frame
                paras = [p for p in tf.paragraphs
                         if "".join(run.text for run in p.runs).strip()]
                if not paras:
                    continue

                # AUTO_SHAPE / CALLOUT / FREEFORM with text: one entry in
                # get_element_boxes if area < 80 % of slide; skip otherwise.
                _is_autoshape_box_r = (
                    _stype_r in (1, 2, 5) and
                    not getattr(shape, "is_placeholder", False)
                )
                if _is_autoshape_box_r:
                    shape_area_r = max(1, r_r - l_r) * max(1, b_r - t_r)
                    if shape_area_r >= slide_area * 0.80:
                        continue   # full-bleed BG rect — get_element_boxes skips it
                    combined = " ".join(
                        "".join(run.text for run in p.runs).strip()
                        for p in paras
                    ).strip()
                    items.append(((t_r + b_r) / 2, combined))
                    continue

                # Regular placeholder / text box: per-paragraph rows.
                # Mirror the row-height and last-row logic of get_element_boxes.
                n_p     = len(paras)
                h_r     = b_r - t_r
                row_h_r = max(4, h_r // n_p)
                for pi, p in enumerate(paras):
                    pt_r = t_r + pi * row_h_r
                    # Last row extends to shape bottom (same as get_element_boxes)
                    pb_r = b_r if pi == n_p - 1 else min(b_r, t_r + (pi + 1) * row_h_r)
                    if pb_r <= pt_r + 2:
                        continue   # too thin — get_element_boxes skips these too
                    text = "".join(run.text for run in p.runs).strip()
                    items.append(((pt_r + pb_r) / 2, text))
                continue

            # ── Table shape ───────────────────────────────────────────────────
            if _stype_r == 19 or getattr(shape, "has_table", False):
                try:
                    tbl   = shape.table
                    n_r   = len(tbl.rows)
                    rh_r  = max(4, (b_r - t_r) // max(n_r, 1))
                    for ri in range(n_r):
                        rt_r = t_r + ri * rh_r
                        rb_r = min(b_r, t_r + (ri + 1) * rh_r)
                        if rb_r <= rt_r + 2:
                            continue
                        row_text = " ".join(
                            cell.text_frame.text
                            for cell in tbl.rows[ri].cells
                            if hasattr(cell, "text_frame")
                        ).strip()
                        items.append(((rt_r + rb_r) / 2, row_text))
                except Exception:
                    items.append(((t_r + b_r) / 2, ""))
                continue

            # ── Picture shapes ────────────────────────────────────────────────
            # get_element_boxes NEVER adds picture shapes (they are always
            # visible from frame 0 or handled via Ken Burns).  Skipping here
            # keeps the element count identical so element_timing aligns 1-to-1.
            try:
                from pptx.enum.shapes import MSO_SHAPE_TYPE
                is_pic = shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            except Exception:
                is_pic = (_stype_r == 13)
            if is_pic:
                continue   # SKIP — must match get_element_boxes behaviour

            # ── Group shapes (type 6) ─────────────────────────────────────────
            # get_element_boxes only adds groups whose area < 80 % of slide.
            if _stype_r == 6:
                shape_area_r = max(1, r_r - l_r) * max(1, b_r - t_r)
                if shape_area_r < slide_area * 0.80:
                    items.append(((t_r + b_r) / 2, ""))
                # else: full-bleed group — get_element_boxes skips it, so do we
                continue

            # ── Line / connector shapes ───────────────────────────────────────
            if _is_line_r:
                items.append(((t_r + b_r) / 2, ""))

        items.sort(key=lambda x: x[0])
        return [text for _, text in items]
    except Exception:
        import traceback
        traceback.print_exc()
        return []


def _get_char_element_boxes(pptx_path: str, slide_idx: int, img_w: int, img_h: int) -> list:
    """
    Like get_element_boxes but splits every text paragraph into individual
    per-character column boxes so that char_overshoot_scale can animate one
    letter at a time.

    For each paragraph, character widths are measured with PIL (proportional
    scaling so the columns always span the full shape width).

    Non-text shapes (images, tables) are returned at their normal granularity
    (same as get_element_boxes).
    """
    try:
        from pptx import Presentation
        from PIL import ImageFont as _IF

        prs = Presentation(pptx_path)
        if slide_idx >= len(prs.slides):
            return []
        slide   = prs.slides[slide_idx]
        sx      = img_w / prs.slide_width
        sy      = img_h / prs.slide_height
        elements = []

        # Ordered list of font paths to try for glyph-width measurement
        _FONT_PATHS = [
            "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
            "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
            "/usr/share/fonts/truetype/freefont/FreeSansBold.ttf",
            "/usr/share/fonts/truetype/ubuntu/Ubuntu-B.ttf",
        ]

        def _load_font(size_px):
            for fp in _FONT_PATHS:
                try:
                    return _IF.truetype(fp, max(8, size_px))
                except Exception:
                    pass
            try:
                return _IF.load_default(max(8, size_px))
            except Exception:
                return _IF.load_default()

        def _char_width(font, ch):
            try:
                if hasattr(font, "getlength"):
                    return max(1.0, float(font.getlength(ch)))
                elif hasattr(font, "getsize"):
                    return max(1.0, float(font.getsize(ch)[0]))
            except Exception:
                pass
            return max(1.0, font.size * 0.6)

        for shape in slide.shapes:
            sl = max(0, int(shape.left  * sx))
            st = max(0, int(shape.top   * sy))
            sr = min(img_w, int((shape.left + shape.width)  * sx))
            sb = min(img_h, int((shape.top  + shape.height) * sy))
            if sr <= sl + 10 or sb <= st + 10:
                continue

            # ── Text shape: split into per-character column boxes ────────────
            if shape.has_text_frame:
                tf    = shape.text_frame
                paras = [p for p in tf.paragraphs
                         if "".join(r.text for r in p.runs).strip()]
                if not paras:
                    continue
                n_p     = len(paras)
                shape_h = sb - st
                row_h   = max(4, shape_h // n_p)

                for pi, para in enumerate(paras):
                    pt = st + pi * row_h
                    pb = min(sb, st + (pi + 1) * row_h)
                    if pb <= pt + 2:
                        continue

                    para_text = "".join(r.text for r in para.runs).strip()
                    if not para_text:
                        continue

                    # Measure glyphs at ~75 % of row height for accuracy
                    font_size_px = max(8, int((pb - pt) * 0.75))
                    font         = _load_font(font_size_px)

                    char_widths = [_char_width(font, ch) for ch in para_text]
                    total_w     = sum(char_widths) or 1.0
                    # Scale so columns fill the full shape width
                    scale_x = (sr - sl) / total_w

                    # Use the full paragraph row height for each character column.
                    # This ensures that when we copy pixels from the rendered slide
                    # image (full_np) we restore exactly the right region — no
                    # vertical clipping of ascenders / descenders.
                    x = float(sl)
                    for ch, cw in zip(para_text, char_widths):
                        cl = int(x)
                        cr = min(sr, int(x + cw * scale_x))
                        if cr > cl:
                            elements.append((cl, pt, cr, pb, "char"))
                        x += cw * scale_x

                    # Close any rounding gap in the last character
                    if elements and elements[-1][4] == "char":
                        last = elements[-1]
                        elements[-1] = (last[0], last[1], sr, last[3], "char")
                continue

            # ── Table: per-row (unchanged from get_element_boxes) ────────────
            shape_type = getattr(shape, "shape_type", None)
            if shape_type == 19 or getattr(shape, "has_table", False):
                try:
                    tbl    = shape.table
                    n_rows = len(tbl.rows)
                    if n_rows > 0:
                        row_h = max(4, (sb - st) // n_rows)
                        for ri in range(n_rows):
                            rt = st + ri * row_h
                            rb = min(sb, st + (ri + 1) * row_h)
                            if rb > rt + 2:
                                elements.append((sl, rt, sr, rb, "table_row"))
                except Exception:
                    elements.append((sl, st, sr, sb, "image"))
                continue

            # ── Pictures ─────────────────────────────────────────────────────
            # SKIP — pictures are never animated (always visible from frame 0).
            # Including them would create an element count mismatch with
            # element_timing (which now excludes pictures to match get_element_boxes).
            try:
                from pptx.enum.shapes import MSO_SHAPE_TYPE
                is_pic = shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            except Exception:
                is_pic = (shape_type == 13)
            if is_pic:
                continue   # SKIP — must match get_element_boxes behaviour

        # Sort top→bottom by row midpoint, left→right within each row
        elements.sort(key=lambda e: ((e[1] + e[3]) / 2, e[0]))
        return elements
    except Exception:
        import traceback
        traceback.print_exc()
        return []


def _sample_bg_color(img_np, l: int, t: int, r: int, b: int) -> tuple:
    """Return the median RGB of pixels just outside the bounding box."""
    import numpy as np
    h, w = img_np.shape[:2]
    pad  = 20
    strips = []
    if t >= pad:
        strips.append(img_np[t - pad : t,         max(0, l) : min(w, r)])
    if b + pad <= h:
        strips.append(img_np[b : min(h, b + pad),  max(0, l) : min(w, r)])
    if l >= pad:
        strips.append(img_np[max(0, t) : min(h, b), l - pad : l])
    if r + pad <= w:
        strips.append(img_np[max(0, t) : min(h, b), r : min(w, r + pad)])

    pixels = [px for s in strips if s.size > 0 for px in s.reshape(-1, 3).tolist()]
    if not pixels:
        return (255, 255, 255)
    med = np.median(pixels, axis=0).astype(int)
    return (int(med[0]), int(med[1]), int(med[2]))


def _get_slide_bg_color(img_np) -> tuple:
    """
    Sample the slide background color robustly.

    Uses four corner regions, but if the top corners are significantly darker
    than the bottom corners (dark header bar), only the bottom corners are used
    so the detected color reflects the main content-area background rather than
    the header band.  This prevents erase fills from picking up the header color
    (e.g. dark blue/black) on slides that have a coloured title bar.
    """
    import numpy as np
    h, w = img_np.shape[:2]
    cs   = min(30, h // 4, w // 4)

    tl = img_np[:cs, :cs].reshape(-1, 3)
    tr = img_np[:cs, w - cs:].reshape(-1, 3)
    bl = img_np[h - cs:, :cs].reshape(-1, 3)
    br = img_np[h - cs:, w - cs:].reshape(-1, 3)

    # Per-corner median brightness (0-255)
    def _bright(px):
        return float(np.median(px)) if px.size else 128.0

    top_bright = (_bright(tl) + _bright(tr)) / 2
    bot_bright  = (_bright(bl) + _bright(br)) / 2

    if top_bright < bot_bright - 40:
        # Dark header → use only bottom corners for content-area BG
        use_corners = [bl, br]
    elif bot_bright < top_bright - 40:
        # Dark footer → use only top corners
        use_corners = [tl, tr]
    else:
        use_corners = [tl, tr, bl, br]

    all_px = np.concatenate([c for c in use_corners if c.size > 0], axis=0)
    if all_px.size == 0:
        return (255, 255, 255)
    med = np.median(all_px, axis=0).astype(int)
    return (int(med[0]), int(med[1]), int(med[2]))


def _ease(t: float) -> float:
    """Smoothstep ease-in-out: starts and ends gently (S-curve)."""
    t = max(0.0, min(1.0, t))
    return t * t * (3.0 - 2.0 * t)


def _apply_elem_anim(frame, full_np, bg_np,
                     l: int, t_b: int, r: int, b: int,
                     anim_type: str, progress: float):
    """
    Blend one element region of *frame* (numpy array, in-place) from the
    background toward the fully-revealed image at the given progress (0→1).

    Extracted so the same pixel logic is shared between the uniform animation
    path and the voiceover-synced path.

    Smoothstep easing is applied to all motion-based animation types so the
    entrance starts and ends gently rather than cutting in at full speed.
    char_overshoot_scale has its own spring easing and is left untouched.
    """
    import numpy as _np
    import math  as _math

    # Apply ease-in-out to linear progress for smooth entrances/exits.
    # char_overshoot_scale handles its own spring-based timing.
    if anim_type != "char_overshoot_scale":
        progress = _ease(progress)

    if anim_type == "text_fade":
        frame[t_b:b, l:r] = (
            bg_np[t_b:b, l:r]   * (1.0 - progress) +
            full_np[t_b:b, l:r] * progress
        ).astype(_np.uint8)

    elif anim_type == "text_slide_up":
        reveal = int(progress * (b - t_b))
        if reveal > 0:
            frame[b - reveal:b, l:r] = full_np[b - reveal:b, l:r]

    elif anim_type == "text_wipe":
        reveal = int(progress * (r - l))
        if reveal > 0:
            frame[t_b:b, l:l + reveal] = full_np[t_b:b, l:l + reveal]

    elif anim_type == "char_overshoot_scale":
        from PIL import Image as _PILImg
        _c, _omega = 8.0, 14.0
        _t = progress * 0.7
        _sc = max(0.001, min(
            1.0 - _math.exp(-_c * _t) * _math.cos(_omega * _t), 2.5
        ))
        _ew, _eh = r - l, b - t_b
        _cx, _cy = (l + r) // 2, (t_b + b) // 2
        _nw = max(1, int(_ew * _sc));  _nh = max(1, int(_eh * _sc))
        _scaled = _np.array(
            _PILImg.fromarray(full_np[t_b:b, l:r]).resize((_nw, _nh), _PILImg.LANCZOS)
        )
        _dl = _cx - _nw // 2;  _dt = _cy - _nh // 2
        _dr = _dl + _nw;       _db = _dt + _nh
        _sl = max(0, -_dl);    _dl = max(0, _dl)
        _st = max(0, -_dt);    _dt = max(0, _dt)
        _dr = min(_dr, frame.shape[1]);  _db = min(_db, frame.shape[0])
        _sr = _sl + (_dr - _dl);        _sb = _st + (_db - _dt)
        if _dr > _dl and _db > _dt and _sr > _sl and _sb > _st:
            frame[_dt:_db, _dl:_dr] = _scaled[_st:_sb, _sl:_sr]


def build_text_animated_clip(
    full_png: str,
    pptx_path: str,
    slide_idx: int,
    duration: float,
    anim_type: str,   # "text_fade" | "text_slide_up" | "text_wipe" | "char_overshoot_scale"
    out_path: str,
    fps: int = 25,
    anim_dur: float = 1.0,
    out_width: int = 1920,
    out_height: int = 1080,
    element_delays: list = None,  # per-element delay in seconds (voiceover-synced)
):
    """
    Build a single-slide MP4 where the text elements animate in (PIL-based).
    Falls back to a plain fade-in on the whole slide if python-pptx / PIL are
    unavailable or no text shapes are found.
    """
    import numpy as np
    from PIL import Image, ImageDraw

    img      = Image.open(full_png).convert("RGB")
    iw, ih   = img.size
    full_np  = np.array(img)

    # Shape-level elements — used for background erasure regardless of mode
    elements = get_element_boxes(pptx_path, slide_idx, iw, ih) if pptx_path else []

    if not elements:
        # Fallback: whole-slide fade-in via FFmpeg at target resolution
        _scale = (f"scale={out_width}:{out_height}:force_original_aspect_ratio=decrease,"
                  f"pad={out_width}:{out_height}:(ow-iw)/2:(oh-ih)/2")
        _cmd = [
            settings.ffmpeg_binary, "-y",
            "-loop", "1", "-t", str(duration), "-i", os.path.abspath(full_png),
            "-vf", f"{_scale},fade=t=in:st=0:d={min(anim_dur, duration * 0.4)}",
            "-c:v", "libx264", "-pix_fmt", "yuv420p", "-r", str(fps), out_path,
        ]
        subprocess.run(_cmd, capture_output=True, timeout=120)
        return

    # anim_elements is always paragraph-level (matches element_timing count).
    # For char_overshoot_scale, we build a separate char-box map so each
    # paragraph's characters animate individually within its timing window
    # WITHOUT breaking the element_timing count check.
    anim_elements = elements

    # ── Build paragraph→character map for char_overshoot_scale ───────────────
    # Keys: paragraph index (into anim_elements)
    # Values: list of character-column boxes within that paragraph
    _para_char_map = {}   # {para_idx: [(l,t,r,b,kind), …]}
    if anim_type == "char_overshoot_scale" and pptx_path:
        _all_chars = _get_char_element_boxes(pptx_path, slide_idx, iw, ih)
        if _all_chars:
            for _pi, _pelem in enumerate(anim_elements):
                _pl, _pt, _pr, _pb, _ = _pelem
                _chars = [
                    c for c in _all_chars
                    if c[1] >= _pt - 4 and c[3] <= _pb + 4 and
                       c[0] >= _pl - 4 and c[2] <= _pr + 4
                ]
                if _chars:
                    _para_char_map[_pi] = _chars

    # ── Create background PNG (all animated elements hidden with sampled BG) ──
    # Always erase at shape-level so the full text area is blanked before
    # animation draws content back in.
    # A small pixel padding ensures text that renders slightly outside its
    # declared bounding box (e.g. descenders, anti-aliased edges) is also
    # covered, preventing "ghost" pixels of later lines showing at frame 0.
    _ERASE_PAD      = 6    # px on each side for text / table rows (increased from 4)
    _LINE_ERASE_PAD = 20   # extra pad for lines — arrowheads extend well past the bbox
    bg_img = img.copy()
    _draw  = ImageDraw.Draw(bg_img)
    # Pre-compute slide background from corners — used for shapes with their own
    # colored backgrounds (LINE, PICTURE) so adjacent-pixel sampling doesn't pick
    # up neighboring shape colors and produce wrong-colored erase fills.
    _slide_bg_color = _get_slide_bg_color(full_np)
    for (l, t, r, b, _kind) in elements:
        _pad = _LINE_ERASE_PAD if _kind == "line" else _ERASE_PAD
        el = max(0, l - _pad)
        et = max(0, t - _pad)
        er = min(iw, r + _pad)
        eb = min(ih, b + _pad)
        if _kind in ("image", "line"):
            # Shapes with their own color (lines, pictures) — use slide BG
            color = _slide_bg_color
        else:
            # Text / table rows — sample just outside the element (more accurate)
            color = _sample_bg_color(full_np, el, et, er, eb)
        _draw.rectangle([el, et, er, eb], fill=color)
    bg_np = np.array(bg_img)

    n_elems   = len(anim_elements)
    frame_dur = round(1.0 / fps, 6)

    frame_dir  = out_path + "_frs"
    concat_txt = out_path + "_frs.txt"
    os.makedirs(frame_dir, exist_ok=True)
    # frame_list: [(path, duration_sec), …]
    frame_list = []
    _fi_ctr    = [0]   # mutable counter shared by inner helper

    def _save_frame(arr):
        fp = os.path.join(frame_dir, f"f{_fi_ctr[0]:07d}.png")
        Image.fromarray(arr).save(fp)
        _fi_ctr[0] += 1
        return fp

    def _compose(revealed_set, active_ei=None, progress=0.0):
        """Build one numpy frame: revealed elements fully shown, active element animating."""
        f = bg_np.copy()
        for ri in revealed_set:
            rl, rt, rr, rb, _rk = anim_elements[ri]
            # Restore the PADDED region (same as what was erased) so text that
            # renders slightly outside its bounding box is fully visible — fixes
            # the persistent "first characters cropped" bug.
            _rp = _LINE_ERASE_PAD if _rk == "line" else _ERASE_PAD
            _rel = max(0, rl - _rp); _ret = max(0, rt - _rp)
            _rer = min(iw, rr + _rp); _reb = min(ih, rb + _rp)
            f[_ret:_reb, _rel:_rer] = full_np[_ret:_reb, _rel:_rer]
        if active_ei is not None:
            al, at, ar, ab, _ = anim_elements[active_ei]
            _apply_elem_anim(f, full_np, bg_np, al, at, ar, ab, anim_type, progress)
        return f

    try:
        # ── Voiceover-synced path ─────────────────────────────────────────────
        # Diagnostic: always log element count and whether we can use voiceover sync.
        if element_delays is not None:
            if len(element_delays) != n_elems:
                print(f"[anim] slide {slide_idx+1}: element count MISMATCH — "
                      f"element_timing has {len(element_delays)} entries but "
                      f"get_element_boxes returned {n_elems} → falling back to uniform. "
                      f"Re-run AI sync to regenerate element_timing.json.")
                # Attempt graceful recovery: truncate / pad to match
                if len(element_delays) > n_elems:
                    element_delays = element_delays[:n_elems]
                    print(f"[anim] slide {slide_idx+1}: truncated element_delays to {n_elems}")
                else:
                    # pad with zeros (show immediately) so no element is hidden
                    element_delays = element_delays + [0.0] * (n_elems - len(element_delays))
                    print(f"[anim] slide {slide_idx+1}: padded element_delays to {n_elems}")
            else:
                print(f"[anim] slide {slide_idx+1}: voiceover-sync active — "
                      f"{n_elems} elements, delays {[round(d,1) for d in element_delays]}, "
                      f"slide dur={duration:.1f}s")
        else:
            print(f"[anim] slide {slide_idx+1}: uniform animation — {n_elems} elements, "
                  f"slide dur={duration:.1f}s (no element_timing)")

        if element_delays is not None and len(element_delays) == n_elems:
            # Each element fades/wipes in at the voiceover-synced moment.
            # For char_overshoot_scale: characters within each element animate
            # sequentially (word-by-word typewriter effect) within a time budget
            # proportional to the element's character count.
            # For other modes: the whole element animates over ELEM_ANIM_DUR.
            ELEM_ANIM_DUR = 0.4   # default entrance time (text_fade / wipe / slide_up)
            # char_overshoot_scale: 3 frames per character, capped at 1.5 s per element
            _CHAR_FRAMES = 3

            max_delay = max(0.0, duration - ELEM_ANIM_DUR)

            # ── Scale delays to fit inside the slide duration ─────────────────
            _max_input = max(element_delays) if element_delays else 0
            if _max_input > max_delay and _max_input > 0:
                _scale = max_delay / _max_input
                safe_delays = [round(max(0.0, d * _scale), 3) for d in element_delays]
                print(f"[anim] slide {slide_idx+1}: max delay {_max_input:.1f}s > "
                      f"slide {duration:.1f}s — rescaling ×{_scale:.3f}")
            else:
                safe_delays = [max(0.0, d) for d in element_delays]

            # ── Safety: topmost element (title / heading) must appear first ───
            if n_elems > 0:
                _top_idx = min(range(n_elems), key=lambda i: anim_elements[i][1])
                safe_delays[_top_idx] = 0.0

            # Sort elements by their voiceover delay
            order    = sorted(range(n_elems), key=lambda i: safe_delays[i])
            revealed = set()
            t_cursor = 0.0   # seconds already written to frame_list

            for ei in order:
                d   = safe_delays[ei]
                gap = max(0.0, d - t_cursor)

                # Static hold frame covering the gap before this element
                if gap > frame_dur * 0.5:
                    fp = _save_frame(_compose(revealed))
                    frame_list.append((fp, round(gap, 6)))
                    t_cursor += round(gap, 6)

                t_cursor = max(t_cursor, d)

                # ── Animation burst for this element ─────────────────────────
                _chars = _para_char_map.get(ei, []) if anim_type == "char_overshoot_scale" else []

                if _chars:
                    # char_overshoot_scale: animate each character column
                    # sequentially — gives the word-by-word typewriter effect.
                    n_chars    = len(_chars)
                    burst_dur  = min(n_chars * _CHAR_FRAMES / fps, 1.5)
                    burst_dur  = max(frame_dur * 2, burst_dur)
                    n_burst    = max(2, round(burst_dur * fps))
                    # Only animate if there's time budget remaining
                    if t_cursor + burst_dur <= duration + frame_dur:
                        for _fi in range(n_burst):
                            _t = _fi / max(n_burst - 1, 1)
                            _fr = _compose(revealed)  # already-revealed elements
                            for _ci, _cbox in enumerate(_chars):
                                _cs = _ci / n_chars
                                _ce = (_ci + 1) / n_chars
                                _cp = min(1.0, max(0.0,
                                    (_t - _cs) / max(_ce - _cs, 0.001)))
                                if _cp <= 0:
                                    continue
                                _cl, _ct, _cr, _cb, _ = _cbox
                                if _t >= _ce:
                                    _fr[_ct:_cb, _cl:_cr] = full_np[_ct:_cb, _cl:_cr]
                                else:
                                    _apply_elem_anim(_fr, full_np, bg_np,
                                                     _cl, _ct, _cr, _cb,
                                                     anim_type, _cp)
                            frame_list.append((_save_frame(_fr), frame_dur))
                        t_cursor += burst_dur
                    # else: element appears instantly
                else:
                    # Standard entrance animation for text_fade / text_wipe /
                    # text_slide_up — or char_overshoot_scale fallback when no
                    # char boxes were found (empty / image element).
                    n_ea = max(2, int(fps * ELEM_ANIM_DUR))
                    if t_cursor + ELEM_ANIM_DUR <= duration + frame_dur:
                        for fi in range(n_ea):
                            prog = fi / max(n_ea - 1, 1)
                            fp   = _save_frame(_compose(revealed, active_ei=ei, progress=prog))
                            frame_list.append((fp, frame_dur))
                        t_cursor += ELEM_ANIM_DUR
                # else: element appears instantly (no animation burst) when
                # there is no remaining time budget — avoids over-long clips.

                revealed.add(ei)

            # Final hold: all elements fully visible for the rest of the slide.
            remaining = duration - t_cursor
            if remaining > frame_dur * 0.5:
                fp = _save_frame(full_np.copy())
                frame_list.append((fp, round(remaining, 6)))

            # Safety: always have at least one frame
            if not frame_list:
                fp = _save_frame(full_np.copy())
                frame_list.append((fp, duration))

            # Correct floating-point drift so the clip's total duration is exactly
            # `duration`.  Even 1-2 ms of cumulative error across many slides will
            # throw transition cut points off by a frame.
            actual_total = sum(d for _, d in frame_list)
            drift = duration - actual_total
            if abs(drift) > 1e-6 and frame_list:
                last_fp, last_dur = frame_list[-1]
                frame_list[-1] = (last_fp, round(max(frame_dur, last_dur + drift), 6))

        # ── Uniform animation path ────────────────────────────────────────────
        else:
            # For char_overshoot_scale: flatten all per-paragraph char boxes
            # into one list so that characters animate sequentially across the
            # whole slide (paragraph by paragraph, char by char within each).
            if anim_type == "char_overshoot_scale" and _para_char_map:
                _all_flat_chars = []
                for _pi in range(n_elems):
                    _all_flat_chars.extend(_para_char_map.get(_pi, [anim_elements[_pi]]))
                _unif_elems = _all_flat_chars
            else:
                _unif_elems = anim_elements

            _ue_count = len(_unif_elems)
            if anim_type == "char_overshoot_scale":
                FRAMES_PER_CHAR = 4   # 4 frames/char = 0.16 s per char at 25 fps
                n_anim     = max(2, min(_ue_count * FRAMES_PER_CHAR, int(fps * 5.0)))
                total_anim = n_anim / fps
            else:
                total_anim = min(anim_dur * _ue_count, 3.0)
                n_anim     = max(2, int(fps * total_anim))
            n_hold = max(1, int(fps * max(0.04, duration - total_anim)))

            # ── Safety: topmost element (title / heading) visible from frame 0 ──
            # Mirror the voiceover-synced path's guarantee: the topmost element
            # (smallest t coordinate) is pre-revealed in the background frame so
            # it appears immediately, even for slides whose title is a coloured
            # autoshape banner (orange / red header bars with icon + text).
            # This fixes "title missing at slide start" on slides where the banner
            # is erased to the slide BG colour and then revealed char-by-char.
            if n_elems > 0:
                _ti = min(range(n_elems), key=lambda i: anim_elements[i][1])
                _tl2, _tt2, _tr2, _tb2, _tk2 = anim_elements[_ti]
                _tp2 = _LINE_ERASE_PAD if _tk2 == "line" else _ERASE_PAD
                _trl = max(0, _tl2 - _tp2); _trt = max(0, _tt2 - _tp2)
                _trr = min(iw, _tr2 + _tp2); _trb = min(ih, _tb2 + _tp2)
                # Write the topmost element into bg_np so every animation frame
                # starts with it already visible — safe because bg_np is never
                # written back to disk, only used as a compositing base.
                bg_np[_trt:_trb, _trl:_trr] = full_np[_trt:_trb, _trl:_trr]
                print(f"[anim] slide {slide_idx+1}: topmost element pre-revealed "
                      f"({_tl2},{_tt2},{_tr2},{_tb2} kind={_tk2})")

            for fi in range(n_anim):
                t_frac = fi / max(n_anim - 1, 1)
                frame  = bg_np.copy()
                # Restore already-fully-revealed elements (with padded region,
                # matching the erase step so no 6px border remains erased)
                for bi, (l, t_b, r, b, _ek) in enumerate(_unif_elems):
                    ei_start = bi / _ue_count
                    ei_end   = (bi + 1) / _ue_count
                    progress = min(1.0, max(0.0,
                        (t_frac - ei_start) / max(ei_end - ei_start, 0.001)))
                    if progress <= 0:
                        continue
                    if t_frac >= ei_end:
                        # Use padded reveal (same extents as erase) so the 6px
                        # border erased around each element is always restored.
                        _rkp = _LINE_ERASE_PAD if _ek == "line" else _ERASE_PAD
                        _rl2 = max(0, l - _rkp);  _rt2 = max(0, t_b - _rkp)
                        _rr2 = min(iw, r + _rkp); _rb2 = min(ih, b + _rkp)
                        frame[_rt2:_rb2, _rl2:_rr2] = full_np[_rt2:_rb2, _rl2:_rr2]
                        continue
                    _apply_elem_anim(frame, full_np, bg_np, l, t_b, r, b,
                                     anim_type, progress)
                fp = _save_frame(frame)
                frame_list.append((fp, frame_dur))

            for fi in range(n_hold):
                fp = _save_frame(full_np.copy())
                frame_list.append((fp, frame_dur))

        # ── Write FFmpeg concat list and encode ───────────────────────────────
        with open(concat_txt, "w", encoding="utf-8") as cf:
            for fp, dur in frame_list:
                cf.write(f"file '{fp.replace(chr(92), '/')}'\n")
                cf.write(f"duration {dur}\n")
            if frame_list:
                cf.write(f"file '{frame_list[-1][0].replace(chr(92), '/')}'\n")

        scale_vf = (f"scale={out_width}:{out_height}:force_original_aspect_ratio=decrease,"
                    f"pad={out_width}:{out_height}:(ow-iw)/2:(oh-ih)/2")
        target_frames = max(1, round(duration * fps))
        cmd = [
            settings.ffmpeg_binary, "-y",
            "-f", "concat", "-safe", "0", "-i", concat_txt,
            "-vf", scale_vf,
            "-c:v", "libx264", "-pix_fmt", "yuv420p", "-r", str(fps),
            "-frames:v", str(target_frames),
            out_path,
        ]
        r = subprocess.run(cmd, capture_output=True, text=True, timeout=300)
        if r.returncode != 0:
            raise RuntimeError(f"FFmpeg text-anim error slide {slide_idx + 1}: {r.stderr[-250:]}")

    finally:
        # Clean up temp frames
        try:
            import shutil as _sh
            _sh.rmtree(frame_dir, ignore_errors=True)
        except Exception:
            pass
        try:
            os.remove(concat_txt)
        except Exception:
            pass


def _concat_clips(clip_paths: list, abs_out: str):
    """Concatenate MP4 clips using the FFmpeg concat demuxer. Cleans up clips afterward."""
    list_file = abs_out + "_cl.txt"
    with open(list_file, "w", encoding="utf-8") as f:
        for p in clip_paths:
            f.write(f"file '{p.replace(chr(92), '/')}'\n")
    cmd = [
        settings.ffmpeg_binary, "-y",
        "-f", "concat", "-safe", "0", "-i", list_file,
        "-c", "copy", abs_out,
    ]
    r = subprocess.run(cmd, capture_output=True, text=True, timeout=600)
    for p in clip_paths:
        try:
            os.remove(p)
        except Exception:
            pass
    try:
        os.remove(list_file)
    except Exception:
        pass
    if r.returncode != 0:
        raise RuntimeError(f"FFmpeg concat error: {r.stderr[-300:]}")


def _get_video_duration(path: str) -> float:
    """Return video duration in seconds using ffprobe."""
    import json
    ffprobe = settings.ffmpeg_binary.replace("ffmpeg", "ffprobe").replace("ffmpeg.exe", "ffprobe.exe")
    r = subprocess.run(
        [ffprobe, "-v", "quiet", "-print_format", "json", "-show_streams", path],
        capture_output=True, text=True, timeout=30,
    )
    try:
        for s in json.loads(r.stdout or "{}").get("streams", []):
            if "duration" in s:
                return float(s["duration"])
    except Exception:
        pass
    return 0.0


# ── Ken Burns (slow zoom / pan) detection & clip building ────────────────────

_KB_EFFECTS = ["zoom_in", "zoom_out", "pan_right", "pan_left", "pan_up", "pan_down"]


def _detect_fullscreen_image_slides(pptx_path: str) -> list:
    """
    Return list[bool] — True for every slide whose picture shapes collectively
    cover ≥ 30 % of the slide area (i.e. it is dominated by a large photo).

    Detects both directly inserted pictures (MSO_SHAPE_TYPE.PICTURE = 13) and
    picture placeholders that have been filled with an image (shape_type 14),
    which python-pptx reports as PLACEHOLDER but still exposes a .image attr.
    Returns an empty list if pptx_path is None / unreadable.
    """
    if not pptx_path:
        return []
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(pptx_path)
        slide_area = prs.slide_width * prs.slide_height
        def _check_shape_for_big_image(shape, slide_area, MSO_SHAPE_TYPE):
            """Return True if shape (or any descendant) is a large picture."""
            try:
                # Recurse into GROUP shapes (type 6)
                if shape.shape_type == 6:
                    for child in shape.shapes:
                        if _check_shape_for_big_image(child, slide_area, MSO_SHAPE_TYPE):
                            return True
                    return False
                # Skip shapes that are too small regardless
                if shape.width * shape.height < slide_area * 0.25:
                    return False
                # Direct picture shapes (Insert > Picture)
                is_pic = shape.shape_type == MSO_SHAPE_TYPE.PICTURE
                if not is_pic:
                    # Picture placeholders filled with an image expose .image
                    try:
                        _ = shape.image
                        is_pic = True
                    except Exception:
                        pass
                return is_pic
            except Exception:
                return False

        results = []
        for idx, slide in enumerate(prs.slides):
            big_image = False
            for shape in slide.shapes:
                if _check_shape_for_big_image(shape, slide_area, MSO_SHAPE_TYPE):
                    big_image = True
                    break
            results.append(big_image)
            if big_image:
                print(f"[kenburns-detect] Slide {idx + 1}: full-screen image detected")
        kb_count = sum(results)
        print(f"[kenburns-detect] {kb_count} full-screen image slide(s) out of {len(results)}")
        return results
    except Exception as e:
        print(f"[kenburns-detect] Detection failed: {e}")
        return []


def _build_kenburns_clip(img_path: str, duration: float, width: int, height: int,
                          out_path: str, fps: int = 25, slide_idx: int = 0) -> None:
    """
    Encode a still image into a video with a slow Ken Burns (zoom / pan) effect.

    Effect cycles deterministically through slide_idx so consecutive slides get
    different effects (zoom_in → zoom_out → pan_right → pan_left → pan_up → pan_down).

    Zoom effects use FFmpeg's ``zoompan`` filter (8 % range, centred).
    Pan effects use scale + animated crop (10 % scale-up) — much faster to encode.
    """
    import math as _math
    effect = _KB_EFFECTS[slide_idx % len(_KB_EFFECTS)]
    # Use round() — same rounding as text-anim clips and _apply_alpha_transitions
    # so all clip durations are consistent and filter trim timestamps never overshoot.
    total_frames = max(25, round(duration * fps))
    w, h = int(width), int(height)

    ZOOM_AMT = 0.08  # 8 % zoom range — subtle but clearly visible

    if effect == "zoom_in":
        step = ZOOM_AMT / total_frames
        vf = (
            f"zoompan=z='min(zoom+{step:.8f},1+{ZOOM_AMT})'"
            f":d={total_frames}:x='iw/2-(iw/zoom/2)':y='ih/2-(ih/zoom/2)'"
            f":s={w}x{h}:fps={fps}"
        )
    elif effect == "zoom_out":
        step = ZOOM_AMT / total_frames
        vf = (
            f"zoompan=z='if(eq(on,1),1+{ZOOM_AMT},max(zoom-{step:.8f},1.0))'"
            f":d={total_frames}:x='iw/2-(iw/zoom/2)':y='ih/2-(ih/zoom/2)'"
            f":s={w}x{h}:fps={fps}"
        )
    else:
        # Pan effects: scale to 110 %, then animate the crop origin
        PAN = 1.10
        sw = round(w * PAN); sw += sw % 2
        sh = round(h * PAN); sh += sh % 2
        dx, dy = sw - w, sh - h
        cx, cy = dx // 2, dy // 2

        if effect == "pan_right":
            cx_expr, cy_expr = f"min({dx}*on/{total_frames},{dx})", str(cy)
        elif effect == "pan_left":
            cx_expr, cy_expr = f"max({dx}*(1-on/{total_frames}),0)", str(cy)
        elif effect == "pan_up":
            cx_expr, cy_expr = str(cx), f"min({dy}*on/{total_frames},{dy})"
        else:  # pan_down
            cx_expr, cy_expr = str(cx), f"max({dy}*(1-on/{total_frames}),0)"

        vf = (
            f"scale={sw}:{sh}:force_original_aspect_ratio=increase,"
            f"crop={w}:{h}:x='{cx_expr}':y='{cy_expr}',"
            f"setsar=1"
        )

    cmd = [
        settings.ffmpeg_binary, "-y",
        "-loop", "1", "-i", img_path,
        "-t", str(round(duration, 6)),
        "-vf", vf,
        "-c:v", "libx264", "-pix_fmt", "yuv420p", "-r", str(fps),
        "-frames:v", str(total_frames),
        out_path,
    ]
    result = subprocess.run(cmd, capture_output=True, text=True, timeout=300)
    if result.returncode != 0:
        raise RuntimeError(f"Ken Burns encode error: {result.stderr[-400:]}")


def _apply_alpha_transitions(
    slide_clips: list,
    transition_clip: str,
    width: str,
    height: str,
    abs_out: str,
    overlap_frames: int = 15,
    fps: int = 25,
    slide_durations: list = None,
) -> str:
    """
    Composite the transition clip at each slide cut using a per-window approach.

    For each slide boundary, a short (~1.2 s) window is extracted, the alpha
    transition is overlaid, and all pieces are reassembled with the FFmpeg
    concat demuxer.  This avoids the 32 KB Windows command-line limit that
    a single filter_complex hits at ~20+ slides and prevents the
    'split=N' memory spike with large N.

    Timeline (overlap = 15 frames @ 25 fps = 0.60 s):
        ┌──── slide A ────┬──── slide B ────┐
                     ↑ T-0.60s  ↑ T+0.60s
        base:  plain_A | window | plain_B | ...
        out:   plain_A | overlay(trans) | plain_B | ...
    """
    import shutil as _shutil
    overlap_secs = overlap_frames / fps
    ffmpeg       = settings.ffmpeg_binary
    w, h         = width, height
    VIDEO_FPS    = fps
    ONE_FRAME    = 1.0 / VIDEO_FPS

    # ── 1. Quantize slide durations to whole frames (same rounding as encoder) ─
    if slide_durations and len(slide_durations) == len(slide_clips):
        clip_durs = list(slide_durations)
    else:
        clip_durs = [_get_video_duration(p) for p in slide_clips]

    frame_counts = [max(1, round(d * VIDEO_FPS)) for d in clip_durs]
    clip_durs    = [fc / VIDEO_FPS for fc in frame_counts]

    # ── 2. Concat slide clips into base video ─────────────────────────────────
    base      = abs_out + "_base.mp4"
    list_file = abs_out + "_base_cl.txt"
    with open(list_file, "w", encoding="utf-8") as f:
        for p in slide_clips:
            f.write(f"file '{p.replace(chr(92), '/')}'\n")
    r = subprocess.run(
        [ffmpeg, "-y", "-f", "concat", "-safe", "0", "-i", list_file,
         "-c", "copy", base],
        capture_output=True, text=True, timeout=600,
    )
    try: os.remove(list_file)
    except Exception: pass
    if r.returncode != 0:
        raise RuntimeError(f"Base concat error: {r.stderr[-200:]}")
    for p in slide_clips:
        try: os.remove(p)
        except Exception: pass

    # ── 3. Compute cut points (frame-accurate, no float drift) ───────────────
    cut_times = []
    frame_acc = 0
    for fc in frame_counts[:-1]:
        frame_acc += fc
        cut_times.append(frame_acc / VIDEO_FPS)

    if not cut_times:
        try: os.rename(base, abs_out)
        except Exception: pass
        return abs_out

    # ── 4. Transition info ────────────────────────────────────────────────────
    trans_dur = _get_video_duration(transition_clip)
    if trans_dur <= 0:
        trans_dur = 2 * overlap_secs
    # Probe actual base duration so tail segment never overshoots
    base_dur_actual = _get_video_duration(base)
    scale_vf = (f"scale={w}:{h}:force_original_aspect_ratio=decrease,"
                f"pad={w}:{h}:(ow-iw)/2:(oh-ih)/2")

    # ── 5. Build segment list ─────────────────────────────────────────────────
    segments  = []   # ("plain"|"overlay", t0, t1)
    prev_end  = 0.0
    for t_cut in cut_times:
        t0 = max(prev_end, t_cut - overlap_secs)
        t1 = min(base_dur_actual, t0 + trans_dur)
        if t1 - t0 < ONE_FRAME:
            continue
        if t0 - prev_end >= ONE_FRAME:
            segments.append(("plain",   prev_end, t0))
        segments.append(    ("overlay", t0,       t1))
        prev_end = t1
    if base_dur_actual - prev_end >= ONE_FRAME:
        segments.append(("plain", prev_end, base_dur_actual))

    # ── 6. Encode each segment individually ──────────────────────────────────
    # Plain segments use -c copy (fast); overlay segments re-encode (~1.2 s each).
    seg_dir   = abs_out + "_segtmp"
    os.makedirs(seg_dir, exist_ok=True)
    seg_paths = []
    timeout_secs = max(1800, int(base_dur_actual * 5))

    for si, (kind, t0, t1) in enumerate(segments):
        seg_path = os.path.join(seg_dir, f"seg{si:05d}.mp4")
        dur      = t1 - t0

        if kind == "plain":
            # Seek-and-copy: fast, uses keyframe alignment.
            # For the first segment (t0=0) use simple trim; for others seek.
            cmd = [
                ffmpeg, "-y",
                "-ss", f"{t0:.6f}", "-i", base,
                "-t",  f"{dur:.6f}",
                "-c", "copy",
                seg_path,
            ]
        else:
            # Overlay the transition clip on this short window
            cmd = [
                ffmpeg, "-y",
                "-ss", f"{t0:.6f}", "-t", f"{dur:.6f}", "-i", base,
                "-i", transition_clip,
                "-filter_complex",
                (f"[0:v]setpts=PTS-STARTPTS[bw];"
                 f"[1:v]{scale_vf}[ts];"
                 f"[bw][ts]overlay=0:0:format=auto:eof_action=pass[out]"),
                "-map", "[out]",
                "-t", f"{dur:.6f}",
                "-c:v", "libx264", "-preset", "fast", "-pix_fmt", "yuv420p",
                seg_path,
            ]

        r = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        if r.returncode != 0:
            # Fallback: copy the plain base segment so the video stays intact
            fb_cmd = [
                ffmpeg, "-y",
                "-ss", f"{t0:.6f}", "-i", base,
                "-t",  f"{dur:.6f}",
                "-c", "copy", seg_path,
            ]
            subprocess.run(fb_cmd, capture_output=True, timeout=60)
        seg_paths.append(seg_path)

    # ── 7. Concat all segments into final output ──────────────────────────────
    _concat_clips(seg_paths, abs_out)

    try: os.remove(base)
    except Exception: pass
    try: _shutil.rmtree(seg_dir, ignore_errors=True)
    except Exception: pass

    return abs_out


def build_video_from_images(
    images: list,
    output_path: str,
    slide_duration: float = 3.0,
    transition: str = "fade",
    resolution: str = "1920x1080",
    durations: list = None,        # per-slide durations; overrides slide_duration when given
    animation: str = "none",       # "none" | "fade_in" | "slide_in_right" | "zoom_in"
                                   # | "text_fade" | "text_slide_up" | "text_wipe"
    pptx_path: str = None,         # needed for text_* animation modes
    transition_clip: str = None,   # path to custom alpha/bumper transition video
    progress_callback=None,        # optional callable(fraction: float 0‥1)
    element_timing: list = None,   # [[delay_sec, …], …] voiceover-synced element delays
):
    """Use FFmpeg to stitch images into a video with transitions."""
    def _report(frac):
        if progress_callback:
            try:
                progress_callback(max(0.0, min(1.0, frac)))
            except Exception:
                pass
    width, height = resolution.split("x")

    # Build per-slide duration list
    if durations is None:
        durations = [slide_duration] * len(images)
    # Pad or trim to match image count
    while len(durations) < len(images):
        durations.append(slide_duration)
    durations = durations[:len(images)]

    # ── Ken Burns detection ───────────────────────────────────────────────────
    # Detect which slides have a dominant full-screen image so we can apply a
    # slow zoom / pan instead of rendering them as a static still.
    kb_slides = _detect_fullscreen_image_slides(pptx_path) if pptx_path else []

    def _make_slide_clip(img: str, dur: float, clip_path: str, slide_i: int) -> None:
        """Build one slide clip: Ken Burns if detected, otherwise plain static."""
        if kb_slides and slide_i < len(kb_slides) and kb_slides[slide_i]:
            effect = _KB_EFFECTS[slide_i % len(_KB_EFFECTS)]
            try:
                _build_kenburns_clip(
                    os.path.abspath(img), dur, int(width), int(height),
                    clip_path, fps=25, slide_idx=slide_i,
                )
                print(f"[kenburns] Slide {slide_i+1}: {effect} applied ✓ ({dur:.1f}s)")
                return
            except Exception as e:
                import traceback as _tb
                print(f"[kenburns] Slide {slide_i+1} fallback to static: {e}")
                _tb.print_exc()
        # Plain static clip
        cmd = [
            settings.ffmpeg_binary, "-y",
            "-loop", "1", "-t", str(dur), "-i", os.path.abspath(img),
            "-vf", (f"scale={width}:{height}:force_original_aspect_ratio=decrease,"
                    f"pad={width}:{height}:(ow-iw)/2:(oh-ih)/2"),
            "-c:v", "libx264", "-pix_fmt", "yuv420p", "-r", "25", clip_path,
        ]
        r = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        if r.returncode != 0:
            raise RuntimeError(f"FFmpeg slide clip error {slide_i+1}: {r.stderr[-200:]}")

    # ── Text animation modes (PIL-based per-element animation) ───────────────
    TEXT_ANIM_MODES = {"text_fade", "text_slide_up", "text_wipe", "char_overshoot_scale"}
    if animation in TEXT_ANIM_MODES:
        clip_paths = []
        abs_out    = os.path.abspath(output_path)
        anim_dur   = 1.0   # seconds for text entrance
        out_w      = int(width)
        out_h      = int(height)

        n_slides = len(images)
        for i, (img, dur) in enumerate(zip(images, durations)):
            _report(i / n_slides)          # progress: 0% … (n-1)/n before concat
            clip = abs_out.replace(".mp4", f"_tclip{i}.mp4")

            # Full-screen image slide → Ken Burns instead of text animation
            if kb_slides and i < len(kb_slides) and kb_slides[i]:
                _kb_effect = _KB_EFFECTS[i % len(_KB_EFFECTS)]
                try:
                    _build_kenburns_clip(
                        os.path.abspath(img), dur, out_w, out_h,
                        clip, fps=25, slide_idx=i,
                    )
                    print(f"[kenburns] Slide {i+1}: {_kb_effect} applied ✓ ({dur:.1f}s)")
                    clip_paths.append(clip)
                    continue
                except Exception as e:
                    import traceback as _tb
                    print(f"[kenburns] Slide {i+1} fallback to text_anim: {e}")
                    _tb.print_exc()

            # Per-slide element delays from voiceover sync (None → uniform fallback)
            slide_elem_delays = (
                element_timing[i]
                if element_timing and i < len(element_timing)
                else None
            )
            try:
                build_text_animated_clip(
                    full_png       = os.path.abspath(img),
                    pptx_path      = pptx_path,
                    slide_idx      = i,
                    duration       = dur,
                    anim_type      = animation,
                    out_path       = clip,
                    fps            = 25,
                    anim_dur       = min(anim_dur, dur * 0.4),
                    out_width      = out_w,
                    out_height     = out_h,
                    element_delays = slide_elem_delays,
                )
            except Exception as e:
                # Text animation failed — log and fall back to plain clip at target res
                import traceback
                print(f"[text_anim] Slide {i+1} failed: {e}")
                traceback.print_exc()
                _make_slide_clip(img, dur, clip, i)
            clip_paths.append(clip)

        _report(0.9)  # 90% — about to concatenate
        if transition_clip and os.path.exists(transition_clip) and len(clip_paths) > 1:
            _apply_alpha_transitions(clip_paths, transition_clip, width, height,
                                     abs_out, slide_durations=durations)
        else:
            _concat_clips(clip_paths, abs_out)
        _report(1.0)
        return

    # ── Per-slide slide/fade/zoom animation ──────────────────────────────────
    if animation != "none":
        anim_dur = 0.8  # seconds for the entrance effect
        clip_paths = []
        abs_out = os.path.abspath(output_path)

        n_slides = len(images)
        for i, (img, dur) in enumerate(zip(images, durations)):
            _report(i / n_slides)
            clip = abs_out.replace(".mp4", f"_aclip{i}.mp4")

            # Full-screen image slide → Ken Burns instead of entrance animation
            if kb_slides and i < len(kb_slides) and kb_slides[i]:
                _kb_eff2 = _KB_EFFECTS[i % len(_KB_EFFECTS)]
                try:
                    _build_kenburns_clip(
                        os.path.abspath(img), dur, int(width), int(height),
                        clip, fps=25, slide_idx=i,
                    )
                    print(f"[kenburns] Slide {i+1}: {_kb_eff2} applied ✓ ({dur:.1f}s)")
                    clip_paths.append(clip)
                    continue
                except Exception as e:
                    import traceback as _tb
                    print(f"[kenburns] Slide {i+1} fallback to anim: {e}")
                    _tb.print_exc()

            if animation == "fade_in":
                vf = (f"fade=t=in:st=0:d={anim_dur},"
                      f"scale={width}:{height}:force_original_aspect_ratio=decrease,"
                      f"pad={width}:{height}:(ow-iw)/2:(oh-ih)/2")
            elif animation == "slide_in_right":
                vf = (f"scale={width}:{height}:force_original_aspect_ratio=decrease,"
                      f"pad={width}:{height}:(ow-iw)/2:(oh-ih)/2,"
                      f"crop={width}:{height}:"
                      f"'if(lt(t,{anim_dur}),(1-t/{anim_dur})*{width},0)':0")
            elif animation == "zoom_in":
                vf = (f"scale={width}:{height}:force_original_aspect_ratio=decrease,"
                      f"pad={width}:{height}:(ow-iw)/2:(oh-ih)/2,"
                      f"zoompan=z='if(lt(t,{anim_dur}),1.1-(0.1*t/{anim_dur}),1)'"
                      f":d=1:x='iw/2-(iw/zoom/2)':y='ih/2-(ih/zoom/2)'"
                      f":s={width}x{height}")
            else:
                vf = (f"scale={width}:{height}:force_original_aspect_ratio=decrease,"
                      f"pad={width}:{height}:(ow-iw)/2:(oh-ih)/2")

            cmd = [
                settings.ffmpeg_binary, "-y",
                "-loop", "1", "-t", str(dur), "-i", os.path.abspath(img),
                "-vf", vf,
                "-c:v", "libx264", "-pix_fmt", "yuv420p",
                "-r", "25",
                clip,
            ]
            r = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            if r.returncode != 0:
                raise RuntimeError(f"FFmpeg animation error slide {i+1}: {r.stderr[-300:]}")
            clip_paths.append(clip)

        _report(0.9)
        if transition_clip and os.path.exists(transition_clip) and len(clip_paths) > 1:
            _apply_alpha_transitions(clip_paths, transition_clip, width, height, abs_out)
        else:
            _concat_clips(clip_paths, abs_out)
        _report(1.0)
        return   # done

    # ── Static (no per-slide animation) ──────────────────────────────────────
    # When a custom transition clip is provided, build individual slide clips
    # then apply alpha transitions between them.
    abs_out = os.path.abspath(output_path)
    if transition_clip and os.path.exists(transition_clip) and len(images) > 1:
        plain_clips = []
        n_slides = len(images)
        for i, (img, dur) in enumerate(zip(images, durations)):
            _report(i / n_slides)
            clip = abs_out.replace(".mp4", f"_plain{i}.mp4")
            _make_slide_clip(img, dur, clip, i)   # Ken Burns if full-screen image detected
            plain_clips.append(clip)
        _report(0.9)
        _apply_alpha_transitions(plain_clips, transition_clip, width, height,
                                 abs_out, slide_durations=durations)
        _report(1.0)
        return

    if transition == "none" or len(images) == 1:
        # concat demuxer — supports per-slide durations natively
        list_file = abs_out + ".txt"
        with open(list_file, "w", encoding="utf-8") as f:
            for img, dur in zip(images, durations):
                abs_img = os.path.abspath(img).replace("\\", "/")
                f.write(f"file '{abs_img}'\n")
                f.write(f"duration {dur}\n")
        cmd = [
            settings.ffmpeg_binary,
            "-y", "-f", "concat", "-safe", "0",
            "-i", list_file,
            "-vf", f"scale={width}:{height}:force_original_aspect_ratio=decrease,"
                   f"pad={width}:{height}:(ow-iw)/2:(oh-ih)/2",
            "-c:v", "libx264", "-pix_fmt", "yuv420p",
            abs_out,
        ]
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=600)
        try:
            os.remove(list_file)
        except Exception:
            pass
    else:
        # xfade transitions — use per-slide durations for each segment
        # Quantize each duration to a 25 fps frame boundary before accumulating
        # offsets; without this, floating-point drift reaches 3-4 frames by the
        # 6-7 minute mark (~80-90 slides at ~5 s/slide).
        _XFADE_FPS = 25
        q_durations = [round(round(d * _XFADE_FPS) / _XFADE_FPS, 6) for d in durations]

        inputs = []
        for img, dur in zip(images, q_durations):
            inputs += ["-loop", "1", "-t", str(dur + 0.5), "-i", img]

        n = len(images)
        filter_parts = []
        for i in range(n):
            filter_parts.append(
                f"[{i}:v]scale={width}:{height}:force_original_aspect_ratio=decrease,"
                f"pad={width}:{height}:(ow-iw)/2:(oh-ih)/2,setsar=1[v{i}]"
            )

        xfade_parts = []
        prev = "v0"
        offset = q_durations[0] - 0.5
        for i in range(1, n):
            out = f"xf{i}"
            xfade_parts.append(
                f"[{prev}][v{i}]xfade=transition={transition}:duration=0.5:offset={offset:.6f}[{out}]"
            )
            prev = out
            offset += q_durations[i] - 0.5

        filter_complex = ";".join(filter_parts + xfade_parts)
        cmd = [settings.ffmpeg_binary, "-y"] + inputs + [
            "-filter_complex", filter_complex,
            "-map", f"[{prev}]",
            "-c:v", "libx264", "-pix_fmt", "yuv420p",
            abs_out,
        ]
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=600)

    if result.returncode != 0:
        raise RuntimeError(f"FFmpeg error: {result.stderr[-500:]}")


def run_slides_to_video(
    job_id: str,
    pptx_path: str,
    slide_duration: float,
    transition: str,
    resolution: str,
    animation: str = "none",
    transition_clip_id: str = None,   # filename inside video_output_dir/transitions/
):
    job = job_store.get(job_id)

    # Persistent frames directory — NOT deleted after conversion so merge can reuse them
    frames_dir = str(Path(settings.video_output_dir) / f"{job_id}_frames")
    Path(frames_dir).mkdir(parents=True, exist_ok=True)

    # Resolve custom transition clip path
    transition_clip_path = None
    if transition_clip_id:
        candidate = str(Path(settings.video_output_dir) / "transitions" / transition_clip_id)
        if os.path.exists(candidate):
            transition_clip_path = candidate

    try:
        job.update(status="processing", progress=10, message="Converting slides to images...")

        # Copy PPTX into frames dir so sync can extract text later
        import shutil as _shutil
        import json as _json
        pptx_copy = str(Path(frames_dir) / "source.pptx")
        _shutil.copy2(pptx_path, pptx_copy)

        images = convert_pptx_to_images(pptx_path, frames_dir)
        if not images:
            raise RuntimeError("No slides found in the uploaded file.")

        # Extract and save per-slide text for AI sync
        slide_data = extract_slide_texts(pptx_copy, len(images))
        with open(str(Path(frames_dir) / "slide_texts.json"), "w", encoding="utf-8") as f:
            _json.dump(slide_data["bodies"], f, ensure_ascii=False)
        with open(str(Path(frames_dir) / "slide_titles.json"), "w", encoding="utf-8") as f:
            _json.dump(slide_data["titles"], f, ensure_ascii=False)

        job.update(progress=50, message=f"Stitching {len(images)} slides into video...")
        output_filename = f"{job_id}_slides.mp4"
        output_path = str(Path(settings.video_output_dir) / output_filename)

        def _stitch_progress(frac, _job=job, _n=len(images)):
            # Maps 0‥1 → 50%‥95% so the bar moves through rendering
            p = 50 + int(frac * 45)
            done = max(1, round(frac * _n))
            _job.update(progress=p, message=f"Rendering slide {done}/{_n}…")

        build_video_from_images(
            images, output_path, slide_duration, transition, resolution,
            animation=animation,
            pptx_path=pptx_copy,
            transition_clip=transition_clip_path,
            progress_callback=_stitch_progress,
        )

        job.update(
            status="done",
            progress=100,
            message="Slides video created successfully!",
            result={
                "video_url": f"/videos/{output_filename}",
                "filename": output_filename,
                "slide_count": len(images),
                "frames_job_id": job_id,
                "transition": transition,
                "resolution": resolution,
                "animation": animation,
                "transition_clip_id": transition_clip_id,
            },
        )
    except Exception as e:
        import traceback
        tb = traceback.format_exc()
        print(f"[slides] ERROR in run_slides_to_video:\n{tb}")
        # Extract app-code frames (exclude Python stdlib subprocess internals only)
        tb_lines = tb.strip().splitlines()
        app_frames = [l.strip() for l in tb_lines
                      if "File \"" in l and "subprocess.py" not in l
                      and "importlib" not in l]
        location = " | ".join(app_frames) if app_frames else " | ".join(tb_lines[-4:])
        job.update(error=f"{str(e)}\n\nAt: {location}")
    finally:
        # Only clean up the uploaded PPTX — keep the frames directory
        if os.path.exists(pptx_path):
            try:
                os.remove(pptx_path)
            except Exception:
                pass


@router.post("/convert")
async def convert_slides(
    background_tasks: BackgroundTasks,
    slide_duration: float = Form(3.0),
    transition: str = Form("fade"),
    resolution: str = Form("1920x1080"),
    animation: str = Form("none"),
    transition_clip_id: str = Form(""),    # ID returned by /upload-transition
    file: UploadFile = File(...),
):
    """
    Convert a .pptx file into a silent video.
    Slide frames are kept for per-slide timing during merge.
    """
    if not file.filename.lower().endswith((".pptx", ".ppt", ".pdf")):
        raise HTTPException(status_code=400, detail="Please upload a .pptx, .ppt, or .pdf file.")

    upload_path = str(Path(settings.upload_dir) / f"{file.filename}")
    with open(upload_path, "wb") as f:
        content = await file.read()
        f.write(content)

    job = job_store.create()
    background_tasks.add_task(
        run_slides_to_video, job.job_id, upload_path, slide_duration, transition, resolution,
        animation, transition_clip_id or None,
    )

    return {"job_id": job.job_id, "status": "pending"}


@router.post("/upload-transition")
async def upload_transition(file: UploadFile = File(...)):
    """
    Accept a custom transition clip (MP4, WebM with alpha, MOV ProRes 4444, etc.)
    and store it in the transitions directory.  Returns a transition_clip_id that
    can be passed back to /convert.
    """
    import uuid as _uuid
    allowed = (".mp4", ".mov", ".avi", ".webm", ".mkv")
    if not file.filename.lower().endswith(allowed):
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported format. Allowed: {', '.join(allowed)}"
        )
    trans_dir = Path(settings.video_output_dir) / "transitions"
    trans_dir.mkdir(parents=True, exist_ok=True)
    ext = os.path.splitext(file.filename)[1].lower()
    clip_id  = f"{_uuid.uuid4()}_trans{ext}"
    out_path = str(trans_dir / clip_id)
    with open(out_path, "wb") as fp:
        fp.write(await file.read())
    return {"transition_clip_id": clip_id, "filename": file.filename}


@router.post("/upload-video")
async def upload_video(file: UploadFile = File(...)):
    """Accept a pre-made video file and save it to the video output directory."""
    import uuid
    allowed = (".mp4", ".mov", ".avi", ".mkv", ".webm")
    if not file.filename.lower().endswith(allowed):
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported video format. Allowed: {', '.join(allowed)}"
        )
    ext = os.path.splitext(file.filename)[1].lower()
    filename = f"{uuid.uuid4()}_slides{ext}"
    output_path = str(Path(settings.video_output_dir) / filename)
    with open(output_path, "wb") as f:
        f.write(await file.read())
    return {
        "video_url": f"/videos/{filename}",
        "filename": filename,
        "slide_count": None,      # unknown for uploaded videos
        "frames_job_id": None,    # no frames — auto_fit/per_slide unavailable
        "transition": "none",
        "resolution": "1920x1080",
    }


@router.post("/refresh-texts/{job_id}")
def refresh_texts(job_id: str):
    """
    Re-extract slide titles and body text from the stored source.pptx WITHOUT
    re-converting the whole presentation.  Call this after upgrading the backend
    when old frames dirs have stale / empty slide_titles.json files.
    """
    import json as _json
    frames_dir = Path(settings.video_output_dir) / f"{job_id}_frames"
    pptx_copy  = frames_dir / "source.pptx"

    if not frames_dir.exists():
        raise HTTPException(status_code=404, detail="Frames directory not found.")
    if not pptx_copy.exists():
        raise HTTPException(status_code=404, detail="source.pptx not found in frames dir.")

    frame_files = sorted(frames_dir.glob("slide_*.png"))
    n_slides    = len(frame_files)
    if n_slides == 0:
        raise HTTPException(status_code=404, detail="No slide frames found.")

    slide_data = extract_slide_texts(str(pptx_copy), n_slides)

    with open(str(frames_dir / "slide_texts.json"),  "w", encoding="utf-8") as f:
        _json.dump(slide_data["bodies"], f, ensure_ascii=False)
    with open(str(frames_dir / "slide_titles.json"), "w", encoding="utf-8") as f:
        _json.dump(slide_data["titles"], f, ensure_ascii=False)

    non_empty = sum(1 for t in slide_data["titles"] if t)
    return {
        "job_id": job_id,
        "slide_count": n_slides,
        "titles_extracted": non_empty,
        "titles": slide_data["titles"],
    }


@router.get("/debug-texts/{job_id}")
def debug_texts(job_id: str, slide: int = 1):
    """
    Diagnostic: return every shape found on a given slide (1-based) with its
    text, shape type, placeholder info and name.  Use this when titles are not
    being extracted to understand the PPTX structure.
    """
    frames_dir = Path(settings.video_output_dir) / f"{job_id}_frames"
    pptx_copy  = frames_dir / "source.pptx"
    if not pptx_copy.exists():
        raise HTTPException(status_code=404, detail="source.pptx not found.")

    try:
        from pptx import Presentation
        prs       = Presentation(str(pptx_copy))
        slide_idx = slide - 1
        if slide_idx >= len(prs.slides):
            raise HTTPException(status_code=404, detail=f"Slide {slide} not found.")

        result = []
        for shape in _iter_shapes(prs.slides[slide_idx].shapes):
            ph     = getattr(shape, "placeholder_format", None)
            text   = ""
            if shape.has_text_frame:
                try:
                    text = shape.text_frame.text.strip()[:120]
                except Exception:
                    text = "<error reading text>"
            result.append({
                "name":       shape.name,
                "shape_type": str(shape.shape_type),
                "has_text":   shape.has_text_frame,
                "text":       text,
                "ph_type":    str(ph.type) if ph else None,
                "ph_idx":     ph.idx if ph else None,
            })

        return {"slide": slide, "total_slides": len(prs.slides), "shapes": result}
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc))


@router.get("/frames/{job_id}")
def get_frames(job_id: str):
    """Return thumbnail URLs for each slide frame stored from a previous conversion."""
    frames_dir = Path(settings.video_output_dir) / f"{job_id}_frames"
    if not frames_dir.exists():
        raise HTTPException(status_code=404, detail="Frames not found for this job.")

    frame_files = sorted(frames_dir.glob("slide_*.png"))
    if not frame_files:
        raise HTTPException(status_code=404, detail="No frame images found.")

    return {
        "frames": [
            {"index": i + 1, "url": f"/videos/{job_id}_frames/{f.name}"}
            for i, f in enumerate(frame_files)
        ]
    }
